# -*- coding: utf-8 -*-
import glob
from pptx import Presentation
from pptx.util import Inches, Emu, Pt

files = glob.glob('D:/edge/ppt/*.pptx')
target = files[0]
print(f"Analyzing: {repr(target)}")

prs = Presentation(target)
print(f'Slide size: {prs.slide_width}x{prs.slide_height}')
print(f'In inches: {prs.slide_width/914400:.2f}x{prs.slide_height/914400:.2f}')
print(f'Slides: {len(prs.slides)}')
print()

for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ===")
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        txt = shape.text[:100].replace('\n',' | ') if hasattr(shape,'text') else ''
        if txt:
            print(f"  [{stype}] {shape.name}")
            print(f"    pos=({shape.left},{shape.top}) size=({shape.width},{shape.height})")
            print(f"    text: {txt}")
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    try:
                        c = str(font.color.rgb) if font.color and font.color.type else 'none'
                    except:
                        c = 'scheme'
                    print(f"    run: \"{run.text[:50]}\" sz={font.size} b={font.bold} c={c}")
    print()
