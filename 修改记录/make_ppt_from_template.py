# -*- coding: utf-8 -*-
"""Use NMU template - clone slides, fill text, preserve design."""
import copy, os
import warnings
warnings.filterwarnings("ignore")
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TM = "D:/edge/ppt/南京医科大学PPT模版（蓝）(1).pptx"
OUT = "D:/codex-workspace/final/修改记录/实验课汇报PPT.pptx"

# ========== SLIDE CLONING ==========
def clone_slide(prs, src_idx):
    """Clone a slide by index, return new slide."""
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    # Remove auto shapes
    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local in ("sp", "pic", "grpSp", "graphicFrame", "cxnSp"):
            spTree.remove(child)
    # Copy source shapes
    for child in src.shapes._spTree:
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local in ("sp", "pic", "grpSp", "graphicFrame", "cxnSp"):
            spTree.append(copy.deepcopy(child))
    return new_slide

def replace_text(shape, old_text, new_text):
    """Replace text in a shape if it contains old_text."""
    if not hasattr(shape, "text_frame") or not shape.has_text_frame:
        return False
    if old_text not in shape.text:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
    return True

def set_text_content(shape, new_text):
    """Set full text, keeping first run formatting."""
    if not hasattr(shape, "text_frame") or ! shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = new_text
        for r in first_para.runs[1:]:
            r.text = ""
    for p in tf.paragraphs[1:]:
        p.clear()

def add_tb(slide, left, top, width, height, text, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
    """Add text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = "Microsoft YaHei"
    if color:
        r.font.color.rgb = color
    return tb

# ========== BUILD PPT ==========
prs = Presentation(TM)
print(f"Template: {len(prs.slides)} slides")

# Save original slide count
orig_count = len(prs.slides)

# ====== SLIDE 1: TITLE ======
# Clone slide 1 (title with background), modify text
s1 = clone_slide(prs, 0)  # Index 0 = Slide 1
# Find and change title text
for sh in s1.shapes:
    if hasattr(sh, "text") and "南京医科大学汇报答辩通用模板" in sh.text:
        set_text_content(sh, "医院信息系统课程实验\n修改与润色汇报")
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(36)
                r.font.bold = True
    if hasattr(sh, "text") and "汇报人" in sh.text and "南小医" in sh.text:
        set_text_content(sh, "汇报人：IMOXXY")
    if hasattr(sh, "text") and "2026" not in sh.text:
        # Check if this is the date placeholder
        if hasattr(sh, "placeholder_format"):
            if sh.placeholder_format.idx == 10:  # Date
                set_text_content(sh, "2026 年 6 月")

print("Slide 1: Title done")

# ====== SLIDE 2: TOC ======
# Slide 9 (index 8) has 目录 CONTENTS with 4 group items
s2 = clone_slide(prs, 8)  # Slide 9
# Change "目录 CONTENTS" text
for sh in s2.shapes:
    if hasattr(sh, "text") and "目录" in sh.text and "CONTENTS" in sh.text:
        replace_text(sh, "目录 CONTENTS", "目 录  CONTENTS")
    # The groups have text in them
    if hasattr(sh, "text") and sh.text.strip():
        t = sh.text.strip()
        if t.startswith("0") or t.startswith("1") or t.startswith("2") or t.startswith("3"):
            pass  # Keep numbers

print("Slide 2: TOC done")

# ====== SLIDE 3: OVERVIEW ======
# Slide 35 (index 34): 一段图文页 - left image, right text
s3 = clone_slide(prs, 34)
for sh in s3.shapes:
    if hasattr(sh, "text"):
        t = sh.text.strip()
        if "一段图文页" in t:
            set_text_content(sh, "项目简介")
        elif "请输入你的正文内容" in t:
            set_text_content(sh, (
                "基于 Spring Boot + Vue.js 的医院信息系统\n"
                "包含挂号、病历、医嘱、缴费、检验等核心业务流程\n"
                "数据库使用 MySQL，共 9 张核心表\n"
                "前后端分离架构"
            ))
        elif "南京医科大学" in t and len(t) < 10:
            set_text_content(sh, "HIS 系统概述")

print("Slide 3: Overview done")

# ====== SLIDES 4-10: EXPERIMENTS 1-7 ======
# Use Slide 40 (index 39): 三段文字页 (3 columns) for experiments
# Each exp covers 2-4 tasks, so 3 columns work well

exp_data = [
    ("实验一：登录与用户信息", ["密码框添加眼睛图标切换显隐", "医生名点击弹窗显示详细信息", "文件: Login.vue / Home.vue"]),
    ("实验二：挂号表单优化", ["挂号类型过滤医生下拉框", "mounted 时自动触发过滤", "文件: PushPatiVisit.vue"]),
    ("实验三：病历管理增强", ["病历打印按钮 + 格式化预览", "诊断名/ICD/拼音首字母搜索", "已提交患者绿色勾标识"]),
    ("实验四：医嘱管理优化", ["表单重置防止数值残留", "按类型行变色(药品绿/检橙...)", "文件: MedicalOrder.vue"]),
    ("实验五：缴费支付优化", ["修复异步执行顺序", "支付成功自动关闭页面", "ispaid 显示是/否"]),
    ("实验六：数据可视化", ["医嘱类型分布柱状图", "卡片点击全屏放大图表", "文件: DataView.vue"]),
    ("实验七：检验检查模块", ["编辑/删除/日期查询", "异常值红/黄高亮", "缴费联动自动创建医嘱"]),
]

for idx, (title, items) in enumerate(exp_data):
    s = clone_slide(prs, 39)  # Slide 40 template
    for sh in s.shapes:
        if hasattr(sh, "text"):
            t = sh.text.strip()
            if "三段文字页" in t:
                set_text_content(sh, title)
            elif "请输入小标题" in t:
                # These are the column title boxes
                found_title = False
                for item in items:
                    if not found_title and item.startswith("文件"):
                        set_text_content(sh, "涉及文件")
                        found_title = True
                if not found_title:
                    pass  # Leave as-is for now
            elif "请在这里输入你的正文内容" in t:
                # These are the column body text boxes
                set_text_content(sh, "\n".join(items))
    print(f"Slide {4+idx}: {title} done")

# ====== SLIDE 11: SUMMARY ======
s11 = clone_slide(prs, 34)  # Use same template as overview
for sh in s11.shapes:
    if hasattr(sh, "text"):
        t = sh.text.strip()
        if "一段图文页" in t:
            set_text_content(sh, "总结与收获")
        elif "请输入你的正文内容" in t:
            set_text_content(sh, (
                "完成了 7 个实验课的代码修改与功能优化\n"
                "覆盖登录、挂号、病历、医嘱、缴费、可视化、检验\n"
                "修复了 BOM 编码、依赖冲突、异步执行等工程问题\n"
                "改进了用户体验：密码显隐、搜索过滤、行变色、异常高亮\n"
                "打通了检验申请→医嘱创建→缴费的完整流程"
            ))
        elif "南京医科大学" in t and len(t) < 10:
            set_text_content(sh, "修改成果统计")

print("Slide 11: Summary done")

# ====== SLIDE 12: THANK YOU ======
s12 = clone_slide(prs, 72)  # Slide 73 (感谢页)
for sh in s12.shapes:
    if hasattr(sh, "text"):
        t = sh.text.strip()
        if "感谢页" in t:
            set_text_content(sh, "感谢聆听")
        elif "特别感谢" in t and "导师" in t:
            set_text_content(sh, "感谢丁海龙老师的悉心指导与课程支持")
        elif "感谢参与" in t:
            set_text_content(sh, "感谢同学们在实验过程中的交流与帮助")
        elif "感谢实验室" in t:
            set_text_content(sh, "通过本次 HIS 系统实验课")
        elif "最后感谢" in t:
            set_text_content(sh, "深刻理解了医院信息系统的业务流程与开发实践")
        elif "谆谆教诲" in t:
            set_text_content(sh, "提升了前后端全栈开发与问题解决能力")

print("Slide 12: Thank you done")

# ====== DELETE ORIGINAL TEMPLATE SLIDES ======
# Remove all original slides from the template
# python-pptx doesn"t have a built-in delete_slide, use XML
print(f"\nRemoving {orig_count} original slides...")

# Get all slide references
sldIdLst = prs.presentation.sldIdLst
sldIdElems = list(sldIdLst)
for sldId in sldIdElems:
    rId = sldId.get("r:id")
    if rId:
        # Delete the slide relationship
        try:
            prs.part.drop_rel(rId)
        except:
            pass
    sldIdLst.remove(sldId)

print(f"All original slides removed. New slide count = {len(prs.slides)}")

# ====== SAVE ======
prs.save(OUT)
print(f"\nSaved to: {OUT}")
