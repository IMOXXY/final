# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Emu, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ============ CONSTANTS ============
SLIDE_W = 12192000
SLIDE_H = 6858000

# Colors
DARK_NAVY = RGBColor(0x0D, 0x23, 0x4B)      # 深藏青
MED_BLUE = RGBColor(0x1A, 0x5C, 0xA8)        # 中蓝
LIGHT_BLUE = RGBColor(0x3B, 0x82, 0xC4)       # 亮蓝
ACCENT_GOLD = RGBColor(0xD4, 0xA5, 0x3C)      # 金色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MID_GRAY = RGBColor(0x66, 0x77, 0x88)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
SOFT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)

# Exp colors
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
TEAL = RGBColor(0x00, 0x96, 0x88)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============ HELPERS ============
def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=12, color=DARK_TEXT, line_spacing=1.5, align=PP_ALIGN.LEFT, bold=False, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return txBox

def add_card(slide, left, top, width, height, title, items, title_color=MED_BLUE, icon_text=None):
    """Add a card with title bar and bullet items"""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, WHITE, RGBColor(0xDD, 0xDD, 0xDD))
    card.shadow.inherit = False
    
    # Title bar (top of card)
    bar = add_rect(slide, left+Emu(10), top+Emu(10), width-Emu(20), Emu(360000), title_color)
    
    # Title text
    title_prefix = f"{icon_text} " if icon_text else ""
    add_text_box(slide, left+Emu(120000), top+Emu(50000), width-Emu(240000), Emu(300000),
                 title_prefix + title, font_size=13, bold=True, color=WHITE)
    
    # Items
    y_offset = Emu(420000)
    for item in items:
        add_text_box(slide, left+Emu(120000), top+y_offset, width-Emu(240000), Emu(220000),
                     f"▪ {item}", font_size=10, color=DARK_TEXT)
        y_offset += Emu(220000)
    return card

# ============ SLIDE 1: TITLE ============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Full background
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_NAVY)

# Decorative abstract shapes
add_rect(slide, 0, 0, SLIDE_W, Emu(4200000), RGBColor(0x0A, 0x1D, 0x3F))
add_rect(slide, 0, Emu(4200000), SLIDE_W, Emu(600000), RGBColor(0x12, 0x2D, 0x5A))

# Gold accent line
add_rect(slide, Emu(4500000), Emu(1400000), Emu(3200000), Emu(40000), ACCENT_GOLD)

# Decorative circle
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(-2000000), Emu(4500000), Emu(6000000), Emu(6000000))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x15, 0x32, 0x65)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(8000000), Emu(-2000000), Emu(5000000), Emu(5000000))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x15, 0x32, 0x65)
circle2.line.fill.background()

# Main title
add_text_box(slide, Emu(800000), Emu(1600000), Emu(10600000), Emu(900000),
             "医院信息系统（HIS）实验课", font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Emu(2500000), Emu(2600000), Emu(7200000), Emu(600000),
             "修改与润色工作汇报", font_size=28, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# Gold accent line 2
add_rect(slide, Emu(4500000), Emu(3300000), Emu(3200000), Emu(30000), ACCENT_GOLD)

# Author info
add_text_box(slide, Emu(2500000), Emu(3600000), Emu(7200000), Emu(400000),
             "医学信息学系", font_size=18, bold=True, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

# Date
add_text_box(slide, Emu(2500000), Emu(4200000), Emu(7200000), Emu(300000),
             "2026 年 6 月", font_size=14, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Footer line - "博学至精 明德至善"
add_text_box(slide, Emu(3500000), Emu(6300000), Emu(5200000), Emu(300000),
             "——  博学至精  明德至善  ——", font_size=11, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ============ SLIDE 2: TOC ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Left panel
add_rect(slide, 0, 0, Emu(4200000), SLIDE_H, DARK_NAVY)

# Section title
add_text_box(slide, Emu(300000), Emu(500000), Emu(3600000), Emu(500000),
             "目 录", font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, Emu(1600000), Emu(1100000), Emu(1000000), Emu(30000), ACCENT_GOLD)

# Left side numbers (1-9)
numbers = ["01", "02", "03", "04", "05", "06", "07", "08", "09"]
for i, n in enumerate(numbers):
    add_text_box(slide, Emu(600000), Emu(1500000 + i*500000), Emu(800000), Emu(400000),
                 n, font_size=16, bold=True, color=RGBColor(0x44, 0x77, 0xBB), align=PP_ALIGN.LEFT)

# Right side items
items = [
    "项目简介与环境配置",
    "实验课 1 — 登录与用户信息",
    "实验课 2 — 挂号表单优化",
    "实验课 3 — 病历管理",
    "实验课 4 — 医嘱管理",
    "实验课 5 — 缴费与支付",
    "实验课 6 — 数据可视化",
    "实验课 7 — 检验检查",
    "总结与收获"
]

for i, item in enumerate(items):
    bg_color = SOFT_BLUE if i % 2 == 0 else WHITE
    y_pos = Emu(500000 + i*650000)
    add_rect(slide, Emu(4600000), y_pos, Emu(7200000), Emu(580000), bg_color)
    add_text_box(slide, Emu(4800000), y_pos+Emu(80000), Emu(6800000), Emu(400000),
                 item, font_size=15, bold=(i==0), color=DARK_TEXT, align=PP_ALIGN.LEFT)

# Gold accent line on left
add_rect(slide, Emu(4200000), Emu(400000), Emu(40000), Emu(6000000), ACCENT_GOLD)

# ============ SLIDE 3: 项目简介与环境配置 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_rect(slide, 0, 0, SLIDE_W, Emu(550000), MED_BLUE)
add_text_box(slide, Emu(400000), Emu(120000), Emu(10000000), Emu(350000),
             "项目简介与环境配置", font_size=22, bold=True, color=WHITE)

# Gold accent under header
add_rect(slide, 0, Emu(550000), SLIDE_W, Emu(30000), ACCENT_GOLD)

# Content - two columns
left_x, right_x = Emu(300000), Emu(6400000)
col_w = Emu(5600000)

# Left card: Project description
add_text_box(slide, left_x, Emu(800000), col_w, Emu(300000),
             "📋 项目概述", font_size=16, bold=True, color=DARK_NAVY)
add_rect(slide, left_x, Emu(1080000), Emu(800000), Emu(30000), MED_BLUE)

desc_lines = [
    "• 基于 Spring Boot + Vue.js 的医院信息系统",
    "• 包含挂号、病历、医嘱、缴费、检验等",
    "      核心医疗业务流程",
    "• 数据库使用 MySQL，共 9 张核心表",
    "• 前后端分离架构",
]
add_multiline_box(slide, left_x, Emu(1200000), col_w, Emu(2000000), desc_lines, font_size=12, line_spacing=1.3)

# Right card: Environment
add_text_box(slide, right_x, Emu(800000), col_w, Emu(300000),
             "⚙️ 环境配置", font_size=16, bold=True, color=DARK_NAVY)
add_rect(slide, right_x, Emu(1080000), Emu(800000), Emu(30000), MED_BLUE)

env_lines = [
    "• JDK 17 (Temurin-17.0.19)",
    "• Node.js 22 LTS + npm",
    "• MySQL 数据库 (his)",
    "• 前后端分离部署",
    "• 修复 BOM 编码与依赖冲突",
]
add_multiline_box(slide, right_x, Emu(1200000), col_w, Emu(2000000), env_lines, font_size=12, line_spacing=1.3)

# Bottom section: Database schema
add_text_box(slide, left_x, Emu(3600000), Emu(11400000), Emu(300000),
             "🗄️  数据库表结构", font_size=16, bold=True, color=DARK_NAVY)
add_rect(slide, left_x, Emu(3880000), Emu(800000), Emu(30000), MED_BLUE)

# Code-like box for DB info
code_bg = add_rounded_rect(slide, left_x, Emu(4000000), Emu(11400000), Emu(2500000), RGBColor(0xF5, 0xF7, 0xFA), RGBColor(0xDD, 0xDD, 0xDD))
db_lines = [
    "tb_user（用户/医生）    tb_patient（患者）     tb_visit（就诊记录）",
    "tb_medical_record（病历）  tb_medical_order（医嘱）  tb_drug_info（药品）",
    "tb_settlement（结算）    tb_check_apply（检验申请）  tb_check_result（检验结果）",
]
add_multiline_box(slide, left_x+Emu(150000), Emu(4100000), Emu(11100000), Emu(2400000), db_lines, font_size=11, color=DARK_TEXT, bold=False, line_spacing=1.8, align=PP_ALIGN.CENTER)

# Footer
add_rect(slide, 0, Emu(6600000), SLIDE_W, Emu(258000), DARK_NAVY)
add_text_box(slide, Emu(400000), Emu(6660000), Emu(11400000), Emu(200000),
             "医院信息系统实验课 · 修改与润色汇报", font_size=9, color=MID_GRAY)

# ============ SLIDES 4-10: Experiments 1-7 ============
experiments = [
    {
        "title": "实验课 1 — 登录与用户信息",
        "left_card": {"title": "密码显示控制", "items": [
            "将两个 v-show 切换的密码框合并为一个",
            "添加可点击眼睛图标切换密码显隐",
            "新增 togglePasswordShow() 方法",
            "文件: src/components/Login.vue",
        ]},
        "right_card": {"title": "用户信息显示", "items": [
            "医生名按钮添加 @click 事件",
            "新增 el-dialog 弹窗展示详细信息",
            "显示用户名、邮箱、医生类型",
            "文件: src/components/Home.vue",
        ]},
        "color": MED_BLUE,
    },
    {
        "title": "实验课 2 — 挂号表单优化",
        "left_card": {"title": "挂号类型过滤", "items": [
            "挂号类型（普通/专家）选择后",
            "自动过滤医生下拉框",
            "mounted 时自动触发过滤",
            "文件: src/components/PushPatiVisit.vue",
        ]},
        "right_card": {"title": "优化效果", "items": [
            "医生选择更精准",
            "减少用户误操作",
            "提升挂号效率",
        ]},
        "color": GREEN,
    },
    {
        "title": "实验课 3 — 病历管理功能增强",
        "left_card": {"title": "病历打印", "items": [
            "新增「打印病历」按钮",
            "el-dialog 弹窗格式化预览",
            "调用浏览器 window.print()",
            "文件: src/components/MedicalRecord.vue",
        ]},
        "right_card": {"title": "搜索与标识", "items": [
            "按诊断名/ICD/拼音首字母搜索",
            "已提交患者显示绿色勾图标",
            "用户体验显著提升",
            "文件: src/components/MedicalRecord.vue",
        ]},
        "color": ORANGE,
    },
    {
        "title": "实验课 4 — 医嘱管理优化",
        "left_card": {"title": "表单重置", "items": [
            "每次打开医嘱弹窗时重置字段",
            "防止上次数值残留",
            "otalOrder 和 orderallprice 清空",
            "文件: src/components/MedicalOrder.vue",
        ]},
        "right_card": {"title": "医嘱行变色", "items": [
            "按 ordertype 显示不同背景色",
            "药品→绿色 / 检查→橙色",
            "治疗→蓝色 / 检验→紫色",
            "文件: src/components/MedicalOrder.vue",
        ]},
        "color": PURPLE,
    },
    {
        "title": "实验课 5 — 缴费与支付优化",
        "left_card": {"title": "支付结果同步", "items": [
            "修复 outPay 异步执行顺序",
            "改用 for...of 逐条处理",
            "支付成功后自动关闭页面",
            "文件: src/components/OutBillPay.vue",
        ]},
        "right_card": {"title": "布尔值显示优化", "items": [
            "ispaid 列显示「是/否」",
            "替代原来的 true/false",
            "更符合中文使用习惯",
            "文件: src/components/MedicalOrder.vue",
        ]},
        "color": TEAL,
    },
    {
        "title": "实验课 6 — 数据可视化增强",
        "left_card": {"title": "第四张卡片图表", "items": [
            "实现「当日医嘱类型分布」柱状图",
            "从 API 获取 code='04' 数据",
            "新增 updateOrderDistChart()",
            "文件: src/components/DataView.vue",
        ]},
        "right_card": {"title": "卡片放大功能", "items": [
            "每张卡片添加全屏图标按钮",
            "点击弹出 el-dialog 大窗口",
            "复制原图表配置到弹窗",
            "文件: src/components/DataView.vue",
        ]},
        "color": RED_ACCENT,
    },
    {
        "title": "实验课 7 — 检验检查模块",
        "left_card": {"title": "编辑/删除/查询/高亮", "items": [
            "新增编辑弹窗与删除确认",
            "新增后端 DELETE 接口",
            "日期范围选择器筛选",
            "异常值红/黄高亮显示",
            "文件: LabCheck.vue + 后端",
        ]},
        "right_card": {"title": "缴费联动", "items": [
            "提交检验申请时自动创建医嘱",
            "ordertype='检验' 类型",
            "接入原有缴费流程",
            "打通检验→缴费全链路",
        ]},
        "color": LIGHT_BLUE,
    },
]

for exp in experiments:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    add_rect(slide, 0, 0, SLIDE_W, Emu(550000), DARK_NAVY)
    add_text_box(slide, Emu(400000), Emu(120000), Emu(10000000), Emu(350000),
                 exp["title"], font_size=20, bold=True, color=WHITE)
    
    # Gold accent line
    add_rect(slide, 0, Emu(550000), SLIDE_W, Emu(30000), ACCENT_GOLD)
    
    # Left card
    lc = exp["left_card"]
    card1 = add_rounded_rect(slide, Emu(300000), Emu(900000), Emu(5600000), Emu(5200000), WHITE, RGBColor(0xDD, 0xDD, 0xDD))
    # Card title bar
    add_rect(slide, Emu(300000)+Emu(10), Emu(900000)+Emu(10), Emu(5600000)-Emu(20), Emu(400000), exp["color"])
    add_text_box(slide, Emu(450000), Emu(960000), Emu(5200000), Emu(300000),
                 lc["title"], font_size=15, bold=True, color=WHITE)
    
    # Left card items
    for j, item in enumerate(lc["items"]):
        is_file = item.startswith("文件:") or item.startswith("新增后端")
        add_text_box(slide, Emu(500000), Emu(1450000 + j*420000), Emu(5200000), Emu(380000),
                     f"✦  {item}", font_size=11, bold=is_file, 
                     color=MID_GRAY if is_file else DARK_TEXT)
    
    # Right card
    rc = exp["right_card"]
    card2 = add_rounded_rect(slide, Emu(6300000), Emu(900000), Emu(5600000), Emu(5200000), WHITE, RGBColor(0xDD, 0xDD, 0xDD))
    # Card title bar
    add_rect(slide, Emu(6300000)+Emu(10), Emu(900000)+Emu(10), Emu(5600000)-Emu(20), Emu(400000), exp["color"])
    add_text_box(slide, Emu(6450000), Emu(960000), Emu(5200000), Emu(300000),
                 rc["title"], font_size=15, bold=True, color=WHITE)
    
    # Right card items
    for j, item in enumerate(rc["items"]):
        is_file = item.startswith("文件:")
        add_text_box(slide, Emu(6500000), Emu(1450000 + j*420000), Emu(5200000), Emu(380000),
                     f"✦  {item}", font_size=11, bold=is_file,
                     color=MID_GRAY if is_file else DARK_TEXT)
    
    # Footer
    add_rect(slide, 0, Emu(6600000), SLIDE_W, Emu(258000), DARK_NAVY)
    add_text_box(slide, Emu(400000), Emu(6660000), Emu(11400000), Emu(200000),
                 "医院信息系统实验课 · 修改与润色汇报", font_size=9, color=MID_GRAY)

# ============ SLIDE 11: Summary ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Full background
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_NAVY)

# Decorative elements
add_rect(slide, 0, Emu(800000), SLIDE_W, Emu(200000), ACCENT_GOLD)

# Main title
add_text_box(slide, Emu(1000000), Emu(1500000), Emu(10200000), Emu(800000),
             "总 结 与 收 获", font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, Emu(4600000), Emu(2400000), Emu(3000000), Emu(30000), ACCENT_GOLD)

# Summary items
summary_items = [
    "✓  完成了 7 个实验课的代码修改与功能优化",
    "✓  覆盖登录、挂号、病历、医嘱、缴费、可视化、检验核心模块",
    "✓  修复了 BOM 编码、依赖冲突、异步执行等工程问题",
    "✓  改进了用户体验：密码显隐、搜索过滤、行变色、异常高亮",
    "✓  打通了检验申请→医嘱创建→缴费的完整业务流程",
]

for i, item in enumerate(summary_items):
    add_text_box(slide, Emu(1500000), Emu(2800000 + i*500000), Emu(9200000), Emu(400000),
                 item, font_size=14, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)
    add_rect(slide, Emu(1200000), Emu(2850000 + i*500000), Emu(50000), Emu(50000), ACCENT_GOLD)

# Modified files count
stat_bg = add_rounded_rect(slide, Emu(3000000), Emu(5800000), Emu(6200000), Emu(700000), RGBColor(0x15, 0x32, 0x65))
add_text_box(slide, Emu(3000000), Emu(5850000), Emu(6200000), Emu(600000),
             "修改涉及 8 个前端组件 + 3 个后端文件  |  共 11 个文件",
             font_size=13, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# Footer line
add_text_box(slide, Emu(3500000), Emu(6500000), Emu(5200000), Emu(300000),
             "——  博学至精  明德至善  ——", font_size=11, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ============ SLIDE 12: Thank you ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_NAVY)

# Decorative circles
c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(-1500000), Emu(-1500000), Emu(5000000), Emu(5000000))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x15, 0x32, 0x65); c1.line.fill.background()

c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(8500000), Emu(3500000), Emu(4000000), Emu(4000000))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x15, 0x32, 0x65); c2.line.fill.background()

add_rect(slide, 0, Emu(3000000), SLIDE_W, Emu(100000), ACCENT_GOLD)

add_text_box(slide, Emu(1000000), Emu(2000000), Emu(10200000), Emu(900000),
             "感谢聆听", font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, Emu(1000000), Emu(3300000), Emu(10200000), Emu(500000),
             "Thank You for Listening", font_size=20, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, Emu(3500000), Emu(6500000), Emu(5200000), Emu(300000),
             "——  博学至精  明德至善  ——", font_size=11, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ============ SAVE ============
output_path = "D:/codex-workspace/final/修改记录/实验课汇报PPT.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
