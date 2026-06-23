# -*- coding: utf-8 -*-
import warnings
warnings.filterwarnings("ignore")
from pptx import Presentation

prs = Presentation("D:/edge/ppt/南京医科大学PPT模版（蓝）(1).pptx")

# Find specific slides
for i, s in enumerate(prs.slides):
    full_text = ""
    for sh in s.shapes:
        if hasattr(sh, "text") and sh.text:
            full_text += sh.text + " "
    
    checks = []
    if "目录" in full_text: checks.append("目录")
    if "感谢" in full_text or "聆听" in full_text: checks.append("感谢/聆听")
    if "Q&A" in full_text.upper() or "QA" in full_text.replace(" ","").upper(): checks.append("Q&A")
    if "总结" in full_text: checks.append("总结")
    if "实验" in full_text: checks.append("实验")
    if "内容" in full_text or "CONTENTS" in full_text.upper(): checks.append("内容/CONTENTS")
    if "南京医科大学汇报答辩通用模板" in full_text: checks.append("TITLE")
    
    if checks:
        print(f"Slide {i+1} (layout={s.slide_layout.name}): {', '.join(checks)}")

# Also look at all layout names and which slides use them
print("\n=== All unique layouts and sample slides ===")
layout_slides = {}
for i, s in enumerate(prs.slides):
    ln = s.slide_layout.name
    if ln not in layout_slides:
        layout_slides[ln] = []
    if len(layout_slides[ln]) < 3:
        layout_slides[ln].append(i+1)

for ln, slides in sorted(layout_slides.items()):
    print(f"  {ln}: slides {slides}")
