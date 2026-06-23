# -*- coding: utf-8 -*-
import glob
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import warnings
warnings.filterwarnings('ignore')

files = glob.glob('D:/edge/ppt/*.pptx')
target = files[0]
prs = Presentation(target)

for si in [0, 4, 5, 11, 12, 13, 26]:
    if si >= len(prs.slides):
        continue
    slide = prs.slides[si]
    print(f"=== Slide {si+1} ===")
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        try:
            txt = shape.text[:80].replace('\n',' | ') if hasattr(shape, 'text') and shape.text else ''
        except:
            txt = ''
        if txt or '矩形' in shape.name or '椭圆' in shape.name:
            print(f"  [{stype}] {shape.name}")
            print(f"    pos=({shape.left},{shape.top}) size=({shape.width},{shape.height})")
        if txt:
            print(f"    text: {txt}")
        try:
            f = shape.fill
            if f.type is not None:
                print(f"    fill_type={f.type}")
                if f.type == 1:
                    try:
                        print(f"    fill_color=#{f.fore_color.rgb}")
                    except:
                        pass
        except:
            pass
    print()
