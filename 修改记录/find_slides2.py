# -*- coding: utf-8 -*-
import warnings
warnings.filterwarnings("ignore")
from pptx import Presentation
from pptx.util import Emu, Pt

prs = Presentation("D:/edge/ppt/南京医科大学PPT模版（蓝）(1).pptx")

# Look at key slides: 1, 4, 9, 35, 40, 44, 48, 51, 61, 73, 114
targets = [0, 3, 8, 34, 39, 43, 47, 50, 60, 72, 113]
for si in targets:
    if si >= len(prs.slides):
        continue
    s = prs.slides[si]
    print(f"=== Slide {si+1} (layout={s.slide_layout.name}) ===")
    for sh in s.shapes:
        stype = str(sh.shape_type)
        try:
            txt = sh.text[:120].replace("\n"," | ") if hasattr(sh, "text") and sh.text else ""
        except:
            txt = ""
        if txt:
            print(f"  [{stype}] {sh.name}: \"{txt}\"")
            print(f"    pos=({sh.left},{sh.top}) size=({sh.width},{sh.height})")
        elif sh.has_text_frame:
            print(f"  [{stype}] {sh.name}: (empty text frame)")
            print(f"    pos=({sh.left},{sh.top}) size=({sh.width},{sh.height})")
        else:
            print(f"  [{stype}] {sh.name}: (no text)")
            print(f"    pos=({sh.left},{sh.top}) size=({sh.width},{sh.height})")
    print()
