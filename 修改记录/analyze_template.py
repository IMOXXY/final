# -*- coding: utf-8 -*-
import glob, os
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.oxml.ns import qn
import warnings
warnings.filterwarnings('ignore')

target = "D:/edge/ppt/南京医科大学PPT模版（蓝）(1).pptx"
print(f"Analyzing: {repr(target)}")

prs = Presentation(target)
print(f"Slide size: {prs.slide_width}x{prs.slide_height}")
print(f"In inches: {prs.slide_width/914400:.2f}x{prs.slide_height/914400:.2f}")
print(f"Slide layouts: {len(prs.slide_layouts)}")
print(f"Slides: {len(prs.slides)}")
print()

# Analyze each layout
for li, layout in enumerate(prs.slide_layouts):
    print(f"=== Layout {li}: {layout.name} ===")
    for ph in layout.placeholders:
        print(f"  Placeholder idx={ph.placeholder_format.idx}, name={ph.name}, type={ph.placeholder_format.type}")
        if ph.has_text_frame:
            txt = ph.text[:50]
            print(f"    text: {txt}")

print()
# Analyze slides
for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ===")
    if slide.slide_layout:
        print(f"  Layout: {slide.slide_layout.name}")
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        try:
            txt = shape.text[:100].replace('\n',' | ') if hasattr(shape, 'text') and shape.text else ''
        except:
            txt = ''
        if txt or '矩形' in shape.name:
            print(f"  [{stype}] {shape.name} pos=({shape.left},{shape.top}) size=({shape.width},{shape.height})")
            if txt: print(f"    text: {txt}")
        try:
            f = shape.fill
            if f.type is not None and f.type != 0:
                print(f"    fill_type={f.type}")
                if f.type == 1:
                    try:
                        print(f"    fill_color=#{f.fore_color.rgb}")
                    except:
                        pass
        except:
            pass
    print()

# Write results to file
import sys
out_path = "D:/codex-workspace/final/修改记录/t_info.txt"
with open(out_path, "w", encoding="utf-8") as outf:
    outf.write(f"Slide size: {prs.slide_width}x{prs.slide_height}\n")
    outf.write(f"Slide count: {len(prs.slides)}\n")
    outf.write(f"Layout count: {len(prs.slide_layouts)}\n\n")
    
    # Analyze first 3 slides more carefully
    for si in [0, 1, 2]:
        if si >= len(prs.slides):
            continue
        slide = prs.slides[si]
        outf.write(f"=== Slide {si+1} === layout={slide.slide_layout.name}\n")
        for shape in slide.shapes:
            try:
                txt = shape.text[:150].replace("\n"," | ") if hasattr(shape, "text") and shape.text else ""
            except:
                txt = ""
            if txt:
                outf.write(f"  Shape {shape.name}: {txt}\n")
                outf.write(f"    pos=({shape.left},{shape.top}) size=({shape.width},{shape.height})\n")
            try:
                f = shape.fill
                if f.type is not None and f.type == 1:
                    try:
                        outf.write(f"    fill=#{f.fore_color.rgb}\n")
                    except:
                        pass
            except:
                pass
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        try:
                            c = str(font.color.rgb) if font.color and font.color.type else "inherit"
                        except:
                            c = "scheme"
                        outf.write(f"    run: \"{run.text[:40]}\" sz={font.size} b={font.bold} c={c}\n")
        outf.write("\n")
    
    # Slide 3 - different layout maybe
    outf.write(f"\n=== Looking at slide layouts ===\n")
    for li, layout in enumerate(prs.slide_layouts):
        if li > 5:
            break
        outf.write(f"  Layout {li}: {layout.name}\n")

print(f"Done, wrote to {out_path}")
